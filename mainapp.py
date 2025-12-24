#  THE UNSTRUCTURED DATA INTEL ENGINE
#  Architecture: Hybrid Streaming + "Data Refinery" Utility
#  Status: PRODUCTION (Fixed: VirtualFile Polymorphism Error)
#
import io
import os
import re
import html
import gc
import time
import csv
import json
import math
import string
import zipfile
import tempfile
import logging
import secrets
import socket
import ipaddress
from dataclasses import dataclass, field
from urllib.parse import urlparse
from collections import Counter
from concurrent.futures import ThreadPoolExecutor
from typing import Dict, List, Tuple, Iterable, Optional, Callable, Any, Union, Set

import numpy as np
import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt
from io import BytesIO
from wordcloud import WordCloud, STOPWORDS
from matplotlib import font_manager
from itertools import pairwise
import openai

# --- graph imports
import networkx as nx
import networkx.algorithms.community as nx_comm
from streamlit_agraph import agraph, Node, Edge, Config

# --- Third Party Imports Checks
try:
    import requests
    from bs4 import BeautifulSoup
except ImportError:
    requests = None
    BeautifulSoup = None

try:
    from scipy.stats import beta as beta_dist
except ImportError:
    beta_dist = None

try:
    from sklearn.decomposition import LatentDirichletAllocation, NMF
    from sklearn.feature_extraction import DictVectorizer
except ImportError:
    LatentDirichletAllocation = None
    NMF = None
    DictVectorizer = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import nltk
    from nltk.sentiment.vader import SentimentIntensityAnalyzer
except ImportError:
    nltk = None
    SentimentIntensityAnalyzer = None

# ==========================================
# ⚙️ CONSTANTS & CONFIGURATION
# ==========================================

MAX_TOPIC_DOCS = 50_000
MAX_SPEAKER_NAME_LENGTH = 30
SENTIMENT_ANALYSIS_TOP_N = 5000
URL_SCRAPE_RATE_LIMIT_SECONDS = 1.0
PROGRESS_UPDATE_MIN_INTERVAL = 100
NPMI_MIN_FREQ = 3
MAX_FILE_SIZE_MB = 200  # Prevent OOM

# Regex Patterns
HTML_TAG_RE = re.compile(r"<[^>]+>")
CHAT_ARTIFACT_RE = re.compile(
    r":\w+:"
    r"|\b(?:monday|tuesday|wednesday|thursday|friday|saturday|sunday|today|yesterday) at \d{1,2}:\d{2}\b"
    r"|\b\d+\s+repl(?:y|ies)\b"
    r"|\d{2}:\d{2}:\d{2}\.\d{3}\s+-->\s+\d{2}:\d{2}:\d{2}\.\d{3}"
    r"|\[[^\]]+\]",
    flags=re.IGNORECASE
)
# Robust URL/Email Regex
URL_EMAIL_RE = re.compile(
    r'(?:https?://|www\.)[a-zA-Z0-9-]+(?:\.[a-zA-Z0-9-]+)+[^\s]*'
    r'|(?:[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,})',
    flags=re.IGNORECASE
)

# Logger Setup
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("IntelEngine")

# Custom Exceptions
class ReaderError(Exception):
    pass

class ValidationError(Exception):
    pass

# ==========================================
# 📦 DATACLASSES
# ==========================================

@dataclass
class CleaningConfig:
    remove_chat: bool = True
    remove_html: bool = True
    remove_urls: bool = True
    unescape: bool = True
    phrase_pattern: Optional[re.Pattern] = None

@dataclass
class ProcessingConfig:
    min_word_len: int = 2
    drop_integers: bool = True
    compute_bigrams: bool = True
    translate_map: Dict[int, Optional[int]] = field(default_factory=dict)
    stopwords: Set[str] = field(default_factory=set)

# ==========================================
# 🛡️ SECURITY & VALIDATION UTILS
# ==========================================

def get_auth_password() -> str:
    pwd = st.secrets.get("auth_password")
    if not pwd:
        st.error("🚨 Configuration Error: 'auth_password' not set in .streamlit/secrets.toml.")
        st.stop()
    return pwd

def validate_url(url: str) -> bool:
    try:
        parsed = urlparse(url)
        if parsed.scheme not in ('http', 'https'):
            return False
        if parsed.hostname in ('localhost', '127.0.0.1', '0.0.0.0', '::1'):
            return False
        return True
    except Exception:
        return False

def validate_sketch_data(data: Dict) -> bool:
    REQUIRED_KEYS = {"total_rows", "counts", "bigrams", "topic_docs"}
    if not isinstance(data, dict): return False
    if not REQUIRED_KEYS.issubset(data.keys()): return False
    if not isinstance(data.get("counts"), dict): return False
    if not isinstance(data.get("topic_docs"), list): return False
    return True

# ==========================================
# 🧠 CORE LOGIC (SCANNER)
# ==========================================

class StreamScanner:
    def __init__(self, doc_batch_size=5):
        self.global_counts = Counter()
        self.global_bigrams = Counter()
        self.total_rows_processed = 0
        self.topic_docs: List[Counter] = []
        self.DOC_BATCH_SIZE = doc_batch_size
        self.limit_reached = False

    def set_batch_size(self, size: int):
        self.DOC_BATCH_SIZE = size

    def update_global_stats(self, counts: Counter, bigrams: Counter, rows: int):
        self.global_counts.update(counts)
        self.global_bigrams.update(bigrams)
        self.total_rows_processed += rows

    def add_topic_sample(self, doc_counts: Counter):
        if not doc_counts: return
        if self.limit_reached: return
        if len(self.topic_docs) >= MAX_TOPIC_DOCS:
            self.limit_reached = True
            return
        self.topic_docs.append(doc_counts)

    def to_json(self) -> str:
        serializable_bigrams = {f"{k[0]}|{k[1]}": v for k, v in self.global_bigrams.items()}
        data = {
            "total_rows": self.total_rows_processed,
            "counts": dict(self.global_counts),
            "bigrams": serializable_bigrams,
            "topic_docs": [dict(c) for c in self.topic_docs],
            "limit_reached": self.limit_reached
        }
        return json.dumps(data)

    def load_from_json(self, json_str: str) -> bool:
        try:
            data = json.loads(json_str)
            if not validate_sketch_data(data):
                return False
            self.total_rows_processed = data.get("total_rows", 0)
            self.global_counts = Counter(data.get("counts", {}))
            raw_bigrams = data.get("bigrams", {})
            self.global_bigrams = Counter()
            for k, v in raw_bigrams.items():
                if "|" in k:
                    parts = k.split("|", 1)
                    self.global_bigrams[(parts[0], parts[1])] = v
            self.topic_docs = [Counter(d) for d in data.get("topic_docs", [])]
            self.limit_reached = data.get("limit_reached", False)
            return True
        except Exception as e:
            logger.error(f"JSON Load Error: {e}")
            return False

# Session State Init
if 'sketch' not in st.session_state: st.session_state['sketch'] = StreamScanner()
if 'total_cost' not in st.session_state: st.session_state['total_cost'] = 0.0
if 'total_tokens' not in st.session_state: st.session_state['total_tokens'] = 0
if 'authenticated' not in st.session_state: st.session_state['authenticated'] = False
if 'auth_error' not in st.session_state: st.session_state['auth_error'] = False
if 'ai_response' not in st.session_state: st.session_state['ai_response'] = ""
if 'last_sketch_hash' not in st.session_state: st.session_state['last_sketch_hash'] = None

def reset_sketch():
    st.session_state['sketch'] = StreamScanner()
    st.session_state['ai_response'] = ""
    st.session_state['last_sketch_hash'] = None
    gc.collect()

def perform_login():
    try:
        correct_password = get_auth_password()
        if secrets.compare_digest(st.session_state.password_input, correct_password):
            st.session_state['authenticated'] = True
            st.session_state['auth_error'] = False
            st.session_state['password_input'] = ""
        else:
            st.session_state['auth_error'] = True
    except Exception:
        st.session_state['auth_error'] = True

def logout():
    st.session_state['authenticated'] = False
    st.session_state['ai_response'] = ""

# ==========================================
# 🛠️ HELPERS & SETUP
# ==========================================

@st.cache_resource(show_spinner="Init NLTK...")
def setup_sentiment_analyzer():
    if nltk is None: return None
    try: nltk.data.find('sentiment/vader_lexicon.zip')
    except LookupError: nltk.download('vader_lexicon')
    return SentimentIntensityAnalyzer()

@st.cache_data(show_spinner=False)
def list_system_fonts() -> Dict[str, str]:
    mapping = {}
    for fe in font_manager.fontManager.ttflist:
        if fe.name not in mapping: mapping[fe.name] = fe.fname
    return dict(sorted(mapping.items(), key=lambda x: x[0].lower()))

def build_punct_translation(keep_hyphens: bool, keep_apostrophes: bool) -> dict:
    punct = string.punctuation
    if keep_hyphens: punct = punct.replace("-", "")
    if keep_apostrophes: punct = punct.replace("'", "")
    return str.maketrans("", "", punct)

def parse_user_stopwords(raw: str) -> Tuple[List[str], List[str]]:
    raw = raw.replace("\n", ",").replace(".", ",")
    phrases, singles = [], []
    for item in [x.strip() for x in raw.split(",") if x.strip()]:
        if " " in item: phrases.append(item.lower())
        else: singles.append(item.lower())
    return phrases, singles

def default_prepositions() -> set:
    return {'about', 'above', 'across', 'after', 'against', 'along', 'among', 'around', 'at', 'before', 'behind', 'below', 'beneath', 'beside', 'between', 'beyond', 'but', 'by', 'concerning', 'despite', 'down', 'during', 'except', 'for', 'from', 'in', 'inside', 'into', 'like', 'near', 'of', 'off', 'on', 'onto', 'out', 'outside', 'over', 'past', 'regarding', 'since', 'through', 'throughout', 'to', 'toward', 'under', 'underneath', 'until', 'up', 'upon', 'with', 'within', 'without'}

def build_phrase_pattern(phrases: List[str]) -> Optional[re.Pattern]:
    if not phrases: return None
    escaped = [re.escape(p) for p in phrases if p]
    if not escaped: return None
    return re.compile(rf"\b(?:{'|'.join(escaped)})\b", flags=re.IGNORECASE)

def estimate_row_count_from_bytes(file_bytes: bytes) -> int:
    if not file_bytes: return 0
    return file_bytes.count(b'\n') + 1

def format_duration(seconds: float) -> str:
    seconds = int(seconds)
    h, r = divmod(seconds, 3600)
    m, s = divmod(r, 60)
    return f"{h:d}:{m:02d}:{s:02d}" if h > 0 else f"{m:d}:{s:02d}"

def make_unique_header(raw_names: List[Optional[str]]) -> List[str]:
    seen: Dict[str, int] = {}
    result: List[str] = []
    for i, nm in enumerate(raw_names):
        name = (str(nm).strip() if nm is not None else "")
        if not name: name = f"col_{i}"
        if name in seen:
            seen[name] += 1
            unique = f"{name}__{seen[name]}"
        else:
            seen[name] = 1
            unique = name
        result.append(unique)
    return result

# --- WEB SCRAPING & VIRTUAL FILES ---
class VirtualFile:
    def __init__(self, name: str, text_content: str):
        self.name = name
        self._bytes = text_content.encode('utf-8')
    
    def getvalue(self) -> bytes:
        return self._bytes
    
    def getbuffer(self) -> memoryview:
        # Added to satisfy the security check relying on .nbytes
        return memoryview(self._bytes)

def fetch_url_content(url: str) -> Optional[str]:
    if not requests or not BeautifulSoup: return None
    
    if not validate_url(url):
        st.toast(f"Blocked unsafe/invalid URL: {url}", icon="🛡️")
        return None
        
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.content, 'html.parser')
        for script in soup(["script", "style", "nav", "footer"]):
            script.decompose()
        return soup.get_text(separator=' ', strip=True)
    except Exception as e:
        st.toast(f"Error fetching {url}: {str(e)}", icon="⚠️")
        return None

# ==========================================
# 📄 FILE READERS
# ==========================================

def read_rows_raw_lines(file_bytes: bytes, encoding_choice: str = "auto") -> Iterable[str]:
    def _iter(enc):
        bio = io.BytesIO(file_bytes)
        with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline=None) as wrapper:
            for line in wrapper: yield line.rstrip("\r\n")
    try:
        if encoding_choice == "latin-1": yield from _iter("latin-1")
        else: yield from _iter("utf-8")
    except UnicodeDecodeError:
        yield "" # Skip bad blocks

def read_rows_vtt(file_bytes: bytes, encoding_choice: str = "auto") -> Iterable[str]:
    for line in read_rows_raw_lines(file_bytes, encoding_choice):
        line = line.strip()
        if not line or line == "WEBVTT" or "-->" in line or line.isdigit(): continue
        if ":" in line:
            parts = line.split(":", 1)
            if len(parts) > 1 and len(parts[0]) < MAX_SPEAKER_NAME_LENGTH and " " in parts[0]:
                yield parts[1].strip()
                continue
        yield line

def read_rows_pdf(file_bytes: bytes) -> Iterable[str]:
    if pypdf is None: 
        st.error("pypdf is missing."); return
    bio = io.BytesIO(file_bytes)
    try:
        reader = pypdf.PdfReader(bio)
        for page in reader.pages:
            text = page.extract_text()
            if text: yield text
    except Exception as e:
        st.warning(f"PDF Read Error: {e}")
        yield ""

def read_rows_pptx(file_bytes: bytes) -> Iterable[str]:
    if pptx is None: return
    bio = io.BytesIO(file_bytes)
    try:
        prs = pptx.Presentation(bio)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    if shape.text: yield shape.text
    except Exception:
        yield ""

def read_rows_json(file_bytes: bytes, selected_key: str = None) -> Iterable[str]:
    bio = io.BytesIO(file_bytes)
    try:
        wrapper = io.TextIOWrapper(bio, encoding="utf-8", errors="replace")
        for line in wrapper:
            if not line.strip(): continue
            try:
                obj = json.loads(line)
                if selected_key and isinstance(obj, dict): yield str(obj.get(selected_key, ""))
                elif isinstance(obj, str): yield obj
                else: yield str(obj)
            except json.JSONDecodeError:
                raise ValueError("Not JSONL") 
    except (ValueError, Exception):
        bio_fallback = io.BytesIO(file_bytes)
        try:
            wrapper = io.TextIOWrapper(bio_fallback, encoding="utf-8", errors="replace")
            data = json.load(wrapper)
            if isinstance(data, list):
                for item in data:
                    if selected_key and isinstance(item, dict): yield str(item.get(selected_key, ""))
                    else: yield str(item)
            elif isinstance(data, dict):
                 if selected_key: yield str(data.get(selected_key, ""))
                 else: yield str(data)
        except json.JSONDecodeError:
            st.warning("Failed to decode JSON file.")

def detect_csv_num_cols(file_bytes: bytes, encoding_choice: str = "auto", delimiter: str = ",") -> int:
    enc = "latin-1" if encoding_choice == "latin-1" else "utf-8"
    bio = io.BytesIO(file_bytes)
    try:
        with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline="") as wrapper:
            rdr = csv.reader(wrapper, delimiter=delimiter)
            row = next(rdr, None)
            return len(row) if row is not None else 0
    except:
        return 0

def get_csv_columns(file_bytes: bytes, encoding_choice: str = "auto", delimiter: str = ",", has_header: bool = True) -> List[str]:
    enc = "latin-1" if encoding_choice == "latin-1" else "utf-8"
    bio = io.BytesIO(file_bytes)
    with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline="") as wrapper:
        rdr = csv.reader(wrapper, delimiter=delimiter)
        first = next(rdr, None)
        if first is None: return []
        return make_unique_header(first) if has_header else [f"col_{i}" for i in range(len(first))]

def get_csv_preview(file_bytes: bytes, encoding_choice: str, delimiter: str, has_header: bool, rows: int = 5) -> pd.DataFrame:
    enc = "latin-1" if encoding_choice == "latin-1" else "utf-8"
    bio = io.BytesIO(file_bytes)
    try:
        df = pd.read_csv(bio, delimiter=delimiter, header=0 if has_header else None, nrows=rows, encoding=enc, on_bad_lines='skip')
        if not has_header: df.columns = [f"col_{i}" for i in range(len(df.columns))]
        return df
    except:
        return pd.DataFrame()

def _iter_tabular_rows(raw_iter, has_header, selected_columns):
    first = next(raw_iter, None)
    if first is None: return

    if has_header:
        header = make_unique_header(list(first))
        name_to_idx = {n: i for i, n in enumerate(header)}
        idxs = [name_to_idx[n] for n in selected_columns if n in name_to_idx]
    else:
        name_to_idx = {f"col_{i}": i for i in range(len(first))}
        idxs = [name_to_idx[n] for n in selected_columns if n in name_to_idx]
        vals = [first[i] if i < len(first) else "" for i in idxs]
        yield vals
    
    for row in raw_iter:
        vals = [row[i] if (row is not None and i < len(row)) else "" for i in idxs]
        yield vals

def iter_csv_selected_columns(file_bytes: bytes, encoding_choice: str, delimiter: str, has_header: bool, selected_columns: List[str], join_with: str) -> Iterable[str]:
    enc = "latin-1" if encoding_choice == "latin-1" else "utf-8"
    bio = io.BytesIO(file_bytes)
    with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline="") as wrapper:
        rdr = csv.reader(wrapper, delimiter=delimiter)
        for vals in _iter_tabular_rows(rdr, has_header, selected_columns):
            vals = [v for v in vals if v]
            yield join_with.join(str(v) for v in vals)

def iter_excel_selected_columns(file_bytes: bytes, sheet_name: str, has_header: bool, selected_columns: List[str], join_with: str) -> Iterable[str]:
    if openpyxl is None: return
    bio = io.BytesIO(file_bytes)
    wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
    ws = wb[sheet_name]
    rows_iter = ws.iter_rows(values_only=True)
    
    for vals in _iter_tabular_rows(rows_iter, has_header, selected_columns):
        vals = [v for v in vals if v]
        yield join_with.join("" if v is None else str(v) for v in vals)
    wb.close()

def get_excel_sheetnames(file_bytes: bytes) -> List[str]:
    if openpyxl is None: return []
    bio = io.BytesIO(file_bytes)
    wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
    sheets = list(wb.sheetnames)
    wb.close()
    return sheets

def get_excel_preview(file_bytes: bytes, sheet_name: str, has_header: bool, rows: int = 5) -> pd.DataFrame:
    if openpyxl is None: return pd.DataFrame()
    bio = io.BytesIO(file_bytes)
    try:
        df = pd.read_excel(bio, sheet_name=sheet_name, header=0 if has_header else None, nrows=rows, engine='openpyxl')
        if not has_header: df.columns = [f"col_{i}" for i in range(len(df.columns))]
        return df
    except:
        return pd.DataFrame()

def excel_estimate_rows(file_bytes: bytes, sheet_name: str, has_header: bool) -> int:
    if openpyxl is None: return 0
    bio = io.BytesIO(file_bytes)
    wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
    ws = wb[sheet_name]
    total = ws.max_row or 0
    wb.close()
    if has_header and total > 0: total -= 1
    return max(total, 0)

# ==========================================
# ⚙️ PROCESSING LOGIC
# ==========================================

def apply_text_cleaning(text: str, config: CleaningConfig) -> str:
    if not isinstance(text, str): 
        return str(text) if text is not None else ""
        
    if config.remove_chat: 
        text = CHAT_ARTIFACT_RE.sub(" ", text)
    if config.remove_html: 
        text = HTML_TAG_RE.sub(" ", text)
    if config.unescape:
        try: text = html.unescape(text)
        except: pass
    if config.remove_urls:
        text = URL_EMAIL_RE.sub(" ", text)
        
    text = text.lower()
    
    if config.phrase_pattern:
        text = config.phrase_pattern.sub(" ", text)
        
    return text.strip()

def process_chunk_iter(
    rows_iter: Iterable[str],
    clean_conf: CleaningConfig,
    proc_conf: ProcessingConfig,
    scanner: StreamScanner,
    progress_cb: Optional[Callable[[int], None]] = None, 
    temp_file_stats: Optional[Counter] = None
):
    _min_len = proc_conf.min_word_len
    _drop_int = proc_conf.drop_integers
    _trans = proc_conf.translate_map
    _stopwords = proc_conf.stopwords
    
    local_global_counts = Counter()
    local_global_bigrams = Counter() if proc_conf.compute_bigrams else Counter()
    
    batch_accum = Counter()
    batch_rows = 0
    is_line_by_line = scanner.DOC_BATCH_SIZE <= 1
    row_count = 0
    update_interval = 2000 

    for line in rows_iter:
        row_count += 1
        text = apply_text_cleaning(line, clean_conf)
        
        filtered_tokens_line: List[str] = []
        for t in text.split():
            t = t.translate(_trans)
            if not t or len(t) < _min_len or (_drop_int and t.isdigit()) or t in _stopwords: continue
            filtered_tokens_line.append(t)
        
        if filtered_tokens_line:
            local_global_counts.update(filtered_tokens_line)
            line_counts = Counter(filtered_tokens_line) 
            
            if proc_conf.compute_bigrams and len(filtered_tokens_line) > 1:
                local_global_bigrams.update(pairwise(filtered_tokens_line))
            
            if is_line_by_line:
                scanner.add_topic_sample(line_counts)
            else:
                batch_accum.update(line_counts)
                batch_rows += 1
                if batch_rows >= scanner.DOC_BATCH_SIZE:
                    scanner.add_topic_sample(batch_accum)
                    batch_accum = Counter()
                    batch_rows = 0

        if progress_cb and (row_count % update_interval == 0): progress_cb(row_count)

    if not is_line_by_line and batch_accum and batch_rows > 0:
        scanner.add_topic_sample(batch_accum)

    scanner.update_global_stats(local_global_counts, local_global_bigrams, row_count)

    if temp_file_stats is not None:
        temp_file_stats.update(local_global_counts)

    if progress_cb: progress_cb(row_count)
    
    del local_global_counts
    del local_global_bigrams
    gc.collect()

def perform_refinery_job(file_obj, chunk_size, clean_conf: CleaningConfig):
    with tempfile.TemporaryDirectory() as temp_dir:
        original_name = os.path.splitext(file_obj.name)[0]
        status_container = st.status(f"⚙️ Refining {file_obj.name}...", expanded=True)
        part_num = 1
        created_files = []
        
        try:
            file_obj.seek(0)
            df_iterator = pd.read_csv(file_obj, chunksize=chunk_size, on_bad_lines='skip', dtype=str)
            
            for chunk in df_iterator:
                for col in chunk.columns:
                    chunk[col] = chunk[col].fillna("")
                    chunk[col] = chunk[col].apply(lambda x: apply_text_cleaning(x, clean_conf))
                
                new_filename = f"{original_name}_cleaned_part_{part_num}.csv"
                temp_path = os.path.join(temp_dir, new_filename)
                chunk.to_csv(temp_path, index=False)
                created_files.append(temp_path)
                status_container.write(f"✅ Processed chunk {part_num} ({len(chunk)} rows)")
                part_num += 1
            
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
                for file_path in created_files:
                    zip_file.write(file_path, arcname=os.path.basename(file_path))
            
            zip_buffer.seek(0)
            status_container.update(label="🎉 Refinery Job Complete!", state="complete", expanded=False)
            return zip_buffer
            
        except Exception as e:
            status_container.update(label="❌ Error", state="error")
            st.error(f"Refinery Error: {str(e)}")
            return None

# ==========================================
# 📊 ANALYTICS
# ==========================================

def render_workflow_guide():
    with st.expander("📘 Comprehensive App Guide: How to use this Tool", expanded=False):
        st.markdown("""
        ### 🌟 What is this?
        This is an **Unstructured Data Intelligence Engine**. It takes "dirty," raw text (from logs, surveys, transcripts, or documents) and extracts mathematical structure, semantic meaning, and qualitative insights without requiring you to write code.

        ---

        ### 🚀 1. Choose Your Workflow

        #### A. The "Direct Scan" (Standard)
        *   **Best for:** PDFs, PowerPoints, Transcripts, or CSVs/JSON files.
        *   **How:** 
            1. Upload your files in the sidebar.
            2. Review the **"Scan Configuration"** box that appears in the main area (select specific columns for CSVs, etc.).
            3. Click the **"⚡ Start Scan"** button for each file; RE-Scan if you've adjusted settings (such as stopwords).
        *   **Result:** The app processes the file into a lightweight statistical "Sketch" and generates a "Quick View" Word Cloud. Once all files are scanned, the aggregate analysis appears below.

        #### B. The "Deep Scan" (Large Datasets)
        *   **Best for:** Massive text dumps (>200MB) where memory is a concern.
        *   **How:** Same as above, but the engine uses **Streaming Mode**.
        *   **Mechanism:** It reads the file in small chunks, extracts statistics, and immediately discards the raw text to save memory.
        *   **Benefit:** This allows you to process datasets larger than your available RAM.

        #### C. The "Enterprise" Workflow (Offline/Secure)
        *   **Best for:** Massive datasets (10M+ rows) or sensitive data (PII) that cannot leave your secure server.
        *   **How:** Use the **Harvester Script** (running locally) to generate a `.json` Sketch file containing only math/counts (no raw text). Upload that JSON here to visualize it.

        ---

        ### 🧠 2. Interpret the Analytics

        #### 🕸️ Network Graph & Community Detection
        *   **The Concept:** Maps how words connect. If "Battery" frequently appears near "Drain," a line connects them.
        *   **Communities (Colors):** Nodes are colored by cluster. If distinct color groups appear, you have successfully separated different conversation topics (e.g., "Login Issues" vs. "Billing Issues").

        #### 🔬 NPMI (Phrase Significance)
        *   **The Problem:** Raw frequency often highlights boring pairs (e.g., "of the", "to be").
        *   **The Solution:** NPMI (Normalized Pointwise Mutual Information) measures *surprise*. It highlights words that appear together *more often than random chance* (e.g., "Credit Card", "Customer Service"). High scores (>0.5) indicate strong semantic links.

        #### 🔍 Bayesian Theme Discovery (Topic Modeling)
        *   **LDA Model:** Best for long text (essays/reports). Assumes documents are a "smoothie" of mixed topics.
        *   **NMF Model:** Best for short text (chats/tickets). Assumes documents fall into distinct, sharp "buckets."

        #### ⚖️ Bayesian Sentiment Inference
        *   **The Value:** Standard sentiment analysis gives you a raw score. This engine calculates a **Credible Interval** (e.g., "We are 95% confident the true positive rate is between 55% and 65%"). This protects you from making business decisions based on small sample sizes.

        ---

        ### ⚡ 3. Utilities
        *   **Data Refinery:** If you have a ~500MB to 1GB CSV that Excel refuses to open, use the Refinery. It cleans the text (removing HTML/Chat logs) and splits it into manageable, Excel-ready chunks.
        *   **AI Analyst:** Uses an LLM (Grok/GPT) to read the *summary statistics* and write a qualitative report.
        """)
def render_use_cases():
    with st.expander("📖 Use-cases", expanded=False):
        st.markdown("""
        ### (Some) use-cases for this unstructured data intelligence engine
        *You'll likely think of more...*

        #### 🏢 Corporate & Strategic
        *   **Customer feedback and support analytics**
        *   **Market and competitive intelligence**
        *   **Brand Voice Audit:** Ensuring consistency in tone across different locations or channels.
        *   **Crisis monitoring in real-time**
        *   **M&A Due Diligence:** Rapidly scanning "Data Rooms" for liabilities (lawsuits, risks) without reading every document.
        *   **Employee engagement analysis**

        #### 📅 Retrospectives & "Year-in-Review"
        *   **The "Narrative Arc":** Ingesting a full year of journals, logs, or chats to visualize how themes shifted from Q1 to Q4.
        *   **Cultural Drift:** Detecting when organizational focus shifted (e.g., from "Innovation" to "Efficiency").
        *   **Tech Debt Monitoring:** visualizing spikes in terms like "quick fix" or "hack" in commit logs over time.

        #### 🔬 Research & Forensics
        *   **Academic and research applications**
        *   **Internal knowledge mining**
        *   **Legal discovery and e-discovery**
        *   **Literary Forensics:** analyzing vocabulary diversity and phrase patterns to detect authorship or ghostwriting.
        *   **Automated discovery of "unknown unknowns"** in large, unstructured datasets.

        #### 🛡️ Security & Privacy
        *   **Compliance and risk monitoring**
        *   **Security:** Insider threat detection.
        *   **The "Privacy Proxy":** Using the *Refinery* tool to strip PII/Chat logs from raw data before sending it to expensive, insecure cloud AI models.

        ---
        #### 🎓 Education: LMS Discussion Forums
        *Insights into what a group of students are discussing:*
        *   Identifying the most common topics and themes students are talking about.
        *   Surfacing frequently asked questions and recurring challenges.
        *   Detecting sentiment trends (e.g., frustration, excitement, confusion) across the class.
        *   Visualizing connections between concepts or issues using network graphs.
        *   Highlighting emerging issues or "unknown unknowns" (e.g., a misunderstood assignment).
        *   Comparing discussion dynamics before and after key events (e.g., exams).
        *   Summarizing participation patterns (who is most/least active).
        *   Providing instructors with actionable summaries for targeted intervention.

        ---
        #### ⚙️ Infrastructure: Middleware for Massive Datasets
        *Acting as an intelligent "translation layer" for 10M+ rows with cost constraints:*
        *   **Statistical Sketches:** Processing raw data into compact sketches (counts, distributions) instead of retaining full text.
        *   **Cost Efficiency:** Massively reducing token usage and API costs by only sending distilled features to LLMs.
        *   **Privacy:** Enabling scalable analysis where no raw text leaves the secure environment.
        *   **Hybrid Workflow:** Supporting iterative, human-in-the-loop workflows (explore summaries -> drill down).
        *   **Pipeline:** Real-time or batch processing for orgs with strict data governance or resource constraints.
        """)
def render_analyst_help():
    with st.expander("🎓 Analyst's Guide", expanded=False):
        st.markdown("""
        **Symptom: Graph vs. Topics Disagree**
        *   **Fix:** Check for 'bridge' words. Lower 'Rows per Document' to 1.
        
        **Symptom: Giant Blob Graph**
        *   **Fix:** Increase 'Min Link Frequency'.
        """)

def calculate_text_stats(counts: Counter, total_rows: int) -> Dict:
    total_tokens = sum(counts.values())
    unique_tokens = len(counts)
    avg_len = sum(len(word) * count for word, count in counts.items()) / total_tokens if total_tokens else 0
    return {
        "Total Rows": total_rows, "Total Tokens": total_tokens,
        "Unique Vocabulary": unique_tokens, "Avg Word Length": round(avg_len, 2),
        "Lexical Diversity": round(unique_tokens / total_tokens, 4) if total_tokens else 0
    }

def calculate_npmi(bigram_counts: Counter, unigram_counts: Counter, total_words: int, min_freq: int = 3) -> pd.DataFrame:
    results = []
    epsilon = 1e-10 
    
    # 1. Safety Check: If no bigrams exist, return empty DataFrame immediately
    if not bigram_counts:
        return pd.DataFrame(columns=["Bigram", "Count", "NPMI"])

    for (w1, w2), freq in bigram_counts.items():
        if freq < min_freq: continue
        count_w1 = unigram_counts.get(w1, 0)
        count_w2 = unigram_counts.get(w2, 0)
        if count_w1 == 0 or count_w2 == 0: continue
        
        prob_w1 = count_w1 / total_words
        prob_w2 = count_w2 / total_words
        prob_bigram = freq / total_words
        
        try: 
            pmi = math.log(prob_bigram / (prob_w1 * prob_w2))
        except ValueError: 
            continue

        log_prob_bigram = math.log(prob_bigram)
        if abs(log_prob_bigram) < epsilon: npmi = 1.0
        else: npmi = pmi / -log_prob_bigram
        results.append({"Bigram": f"{w1} {w2}", "Count": freq, "NPMI": round(npmi, 3)})
    
    # 2. Safety Check: If results list is empty, return empty DataFrame with columns
    df = pd.DataFrame(results)
    if df.empty:
        return pd.DataFrame(columns=["Bigram", "Count", "NPMI"])
        
    return df.sort_values("NPMI", ascending=False)

def perform_topic_modeling(synthetic_docs: List[Counter], n_topics: int, model_type: str) -> Optional[List[Dict]]:
    if not DictVectorizer or len(synthetic_docs) < 1: return None
    vectorizer = DictVectorizer(sparse=True)
    dtm = vectorizer.fit_transform(synthetic_docs)
    n_samples, n_features = dtm.shape
    if n_samples == 0 or n_features == 0: return None
    
    safe_n_topics = min(n_topics, min(n_samples, n_features)) if model_type == "NMF" else min(n_topics, n_samples)
    if safe_n_topics < 1: return None

    model = None
    try:
        if model_type == "LDA": model = LatentDirichletAllocation(n_components=safe_n_topics, random_state=42, max_iter=10)
        elif model_type == "NMF": model = NMF(n_components=safe_n_topics, random_state=42, init='nndsvd')
        model.fit(dtm)
    except ValueError: return None
    
    feature_names = vectorizer.get_feature_names_out()
    topics = []
    for topic_idx, topic in enumerate(model.components_):
        top_indices = topic.argsort()[:-7:-1]
        top_words = [feature_names[i] for i in top_indices]
        strength = sum(topic[i] for i in top_indices)
        topics.append({"id": topic_idx + 1, "words": top_words, "strength": strength})
    return topics

def perform_bayesian_sentiment_analysis(counts: Counter, sentiments: Dict[str, float], pos_thresh: float, neg_thresh: float) -> Optional[Dict]:
    if not beta_dist: return None
    pos_count = sum(counts[w] for w, s in sentiments.items() if s >= pos_thresh)
    neg_count = sum(counts[w] for w, s in sentiments.items() if s <= neg_thresh)
    total_informative = pos_count + neg_count
    if total_informative < 1: return None

    alpha_post = 1 + pos_count
    beta_post = 1 + neg_count
    mean_prob = alpha_post / (alpha_post + beta_post)
    lower_ci, upper_ci = beta_dist.ppf([0.025, 0.975], alpha_post, beta_post)
    x = np.linspace(0, 1, 300)
    y = beta_dist.pdf(x, alpha_post, beta_post)
    return {
        "pos_count": pos_count, "neg_count": neg_count, "total": total_informative,
        "mean_prob": mean_prob, "ci_low": lower_ci, "ci_high": upper_ci,
        "x_axis": x, "pdf_y": y
    }

# ==========================================
# 🎨 VISUALIZATION UTILS
# ==========================================

@st.cache_data(show_spinner="Analyzing term sentiment...")
def get_sentiments(_analyzer, terms: Tuple[str, ...]) -> Dict[str, float]:
    if not _analyzer or not terms: return {}
    return {term: _analyzer.polarity_scores(term)['compound'] for term in terms}

def create_sentiment_color_func(sentiments: Dict[str, float], pos_color, neg_color, neu_color, pos_thresh, neg_thresh):
    def color_func(word, font_size, position, orientation, random_state=None, **kwargs):
        score = sentiments.get(word, 0.0)
        if score >= pos_thresh: return pos_color
        elif score <= neg_thresh: return neg_color
        else: return neu_color
    return color_func

# --- RESTORED HELPER FUNCTION ---
def get_sentiment_category(score: float, pos_threshold: float, neg_threshold: float) -> str:
    if score >= pos_threshold: return "Positive"
    if score <= neg_threshold: return "Negative"
    return "Neutral"

def build_wordcloud_figure_from_counts(counts: Counter, max_words: int, width: int, height: int, bg_color: str, colormap: str, font_path: Optional[str], random_state: int, color_func: Optional[Callable] = None):
    limited = dict(counts.most_common(max_words))
    if not limited: return plt.figure(), None
    wc = WordCloud(width=width, height=height, background_color=bg_color, colormap=colormap, font_path=font_path, random_state=random_state, color_func=color_func, collocations=False, normalize_plurals=False).generate_from_frequencies(limited)
    fig_w, fig_h = max(6.0, width / 100.0), max(3.0, height / 100.0)
    fig, ax = plt.subplots(figsize=(fig_w, fig_h), dpi=100)
    ax.imshow(wc, interpolation="bilinear"); ax.axis("off"); plt.tight_layout()
    return fig, wc

def fig_to_png_bytes(fig: plt.Figure) -> BytesIO:
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", pad_inches=0.1)
    buf.seek(0)
    return buf

# ==========================================
# 🤖 AI LOGIC (Cost Tracking + Chat)
# ==========================================

def call_llm_and_track_cost(system_prompt: str, user_prompt: str, config: dict):
    """
    Generic LLM caller that handles API calls, calculates cost based on tokens,
    and updates the session state accumulators.
    """
    try:
        client = openai.OpenAI(api_key=config['api_key'], base_url=config['base_url'])
        response = client.chat.completions.create(
            model=config['model_name'],
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ]
        )
        
        # Calculate Cost
        in_tok = 0
        out_tok = 0
        
        # Standard OpenAI Usage
        if hasattr(response, 'usage') and response.usage:
            in_tok = response.usage.prompt_tokens
            out_tok = response.usage.completion_tokens
        
        # Cost calculation (Price per 1M tokens)
        cost = (in_tok * config['price_in'] / 1_000_000) + (out_tok * config['price_out'] / 1_000_000)
        
        # Update Session State
        st.session_state['total_tokens'] += (in_tok + out_tok)
        st.session_state['total_cost'] += cost
            
        return response.choices[0].message.content
        
    except Exception as e:
        return f"AI Error: {str(e)}"

# ==========================================
# 🚀 MAIN APP UI
# ==========================================

st.set_page_config(page_title="Intel Engine", layout="wide")
st.title("🧠 Intel Engine: Unstructured Data Analytics")

render_workflow_guide() # calling the updated guide
render_use_cases()      # calling the added use-cases
analyzer = setup_sentiment_analyzer()

# --- SIDEBAR ---
with st.sidebar:
    st.header("📂 Data Input")
    uploaded_files = st.file_uploader("Upload Files", type=["csv", "xlsx", "vtt", "txt", "json", "pdf", "pptx"], accept_multiple_files=True)
    clear_on_scan = st.checkbox("Clear previous data", value=False)
    if st.button("🗑️ Reset All"): reset_sketch(); st.rerun()
    
    st.divider()
    with st.expander("🌐 Web/Text Import"):
        sketch_upload = st.file_uploader("Import Sketch (.json)", type=["json"])
        
        # LOGIC FIX: Hash check for Sketch Upload
        if sketch_upload:
            file_hash = hash(sketch_upload.getvalue())
            if st.session_state.get('last_sketch_hash') != file_hash:
                if st.session_state['sketch'].load_from_json(sketch_upload.getvalue().decode('utf-8')):
                    st.session_state['last_sketch_hash'] = file_hash
                    st.success("Sketch Loaded Successfully!")
                else:
                    st.error("Invalid Sketch File")

        url_input = st.text_area("URLs (one per line)")
        manual_input = st.text_area("Manual Text")
        
    st.divider()
    st.header("🔐 AI Setup")
    if st.session_state['authenticated']:
        st.success("Unlocked")
        
        with st.expander("🤖 Provider Settings", expanded=True):
            ai_provider = st.radio("Provider", ["xAI (Grok)", "OpenAI (GPT-4o)"])
            
            if "OpenAI" in ai_provider:
                api_key_name = "openai_api_key"
                base_url = None 
                model_name = st.selectbox("Model", ["gpt-4o", "gpt-4o-mini"])
                if "mini" in model_name:
                    price_in, price_out = 0.15, 0.60
                else:
                    price_in, price_out = 2.50, 10.00
            else:
                api_key_name = "xai_api_key"
                base_url = "https://api.x.ai/v1"
                model_options = {
                    "Grok 4.1 Fast (Reasoning) [Best Value]": "grok-4-1-fast-reasoning",
                    "Grok 4": "grok-4-0709",
                    "Grok 2 (Legacy)": "grok-2-1212"
                }
                choice = st.selectbox("Model", list(model_options.keys()))
                model_name = model_options[choice]
                
                if "fast" in model_name:
                    price_in, price_out = 0.20, 0.50
                elif "grok-4" in model_name:
                    price_in, price_out = 3.00, 15.00
                else:
                    price_in, price_out = 2.00, 10.00

            api_key = st.secrets.get(api_key_name)
            if not api_key: api_key = st.text_input(f"Enter {api_key_name}", type="password")
            
            ai_config = {
                'api_key': api_key,
                'base_url': base_url,
                'model_name': model_name,
                'price_in': price_in,
                'price_out': price_out
            }

        with st.expander("💰 Cost Estimator", expanded=False):
            c1, c2 = st.columns(2)
            c1.markdown(f"**Tokens:**\n{st.session_state['total_tokens']:,}")
            c2.markdown(f"**Cost:**\n`${st.session_state['total_cost']:.5f}`")
            if st.button("Reset Cost"):
                st.session_state['total_cost'] = 0.0
                st.session_state['total_tokens'] = 0
                st.rerun()

        if st.button("Logout"): logout(); st.rerun()
    else:
        st.text_input("Password", type="password", key="password_input", on_change=perform_login)
        if st.session_state['auth_error']: st.error("Incorrect password")

    st.divider()
    st.header("⚙️ Configuration")
    
    st.markdown("### 🧹 Cleaning")
    clean_conf = CleaningConfig(
        remove_chat=st.checkbox("Remove Chat Artifacts", True),
        remove_html=st.checkbox("Remove HTML", True),
        remove_urls=st.checkbox("Remove URLs", True),
        unescape=st.checkbox("Unescape HTML", True)
    )
    
    st.markdown("### 🛑 Stopwords")
    user_sw = st.text_area("Stopwords (comma-separated)", "firstname.lastname, jane doe")
    phrases, singles = parse_user_stopwords(user_sw)
    clean_conf.phrase_pattern = build_phrase_pattern(phrases)
    stopwords = set(STOPWORDS).union(singles)
    if st.checkbox("Remove Prepositions", True): stopwords.update(default_prepositions())
    
    st.markdown("### ⚙️ Processing")
    proc_conf = ProcessingConfig(
        min_word_len=st.slider("Min Word Len", 1, 10, 2, help="Ignore words shorter than this (e.g., set to 3 to skip 'is', 'it', 'at')."),
        drop_integers=st.checkbox("Drop Integers", True, help="Remove numbers (e.g., IDs, Years). Uncheck if you are analyzing numeric data."),
        compute_bigrams=st.checkbox("Bigrams", True, help="Count pairs of words (e.g., 'customer service') as distinct units."),
        translate_map=build_punct_translation(st.checkbox("Keep Hyphens"), st.checkbox("Keep Apostrophes")),
        stopwords=stopwords
    )
    
    st.markdown("### 🎨 Appearance")
    bg_color = st.color_picker("background color", value="#ffffff")
    colormap = st.selectbox("colormap", options=["viridis", "plasma", "inferno", "magma", "cividis", "tab10", "tab20", "Dark2", "Set3", "rainbow", "cubehelix", "prism", "Blues", "Greens", "Oranges", "Reds", "Purples", "Greys"], index=0)
    max_words = st.slider("max words", 50, 3000, 1000, 50)
    width = st.slider("image width", 600, 2400, 1200, 100)
    height = st.slider("image height", 300, 1400, 600, 50)
    random_state = st.number_input("random seed", 0, value=42, step=1, help="Fix this number to ensure the Word Cloud and Topics look exactly the same every time you run the analysis.")
    
    font_map, font_names = list_system_fonts(), list(list_system_fonts().keys())
    combined_font_name = st.selectbox("font", font_names or ["(default)"], 0)
    combined_font_path = font_map.get(combined_font_name) if font_names else None

    st.markdown("### 📊 Tables")
    top_n = st.number_input("Top Terms to Display", min_value=5, max_value=1000, value=20, help="Controls how many rows appear in the frequency tables below.")

    st.markdown("### 🔬 Sentiment")
    enable_sentiment = st.checkbox("Enable Sentiment", False)
    if enable_sentiment and analyzer is None:
        st.error("NLTK not found.")
        enable_sentiment = False
    
    pos_threshold, neg_threshold, pos_color, neu_color, neg_color = 0.05, -0.05, '#2ca02c', '#808080', '#d62728'
    if enable_sentiment:
        c1, c2 = st.columns(2)
        with c1: pos_threshold = st.slider("pos threshold", 0.0, 1.0, 0.05, 0.01)
        with c2: neg_threshold = st.slider("neg threshold", -1.0, 0.0, -0.05, 0.01)
        c1, c2, c3 = st.columns(3)
        with c1: pos_color = st.color_picker("pos color", value=pos_color)
        with c2: neu_color = st.color_picker("neu color", value=neu_color)
        with c3: neg_color = st.color_picker("neg color", value=neg_color)

    doc_granularity = st.select_slider("Rows per Doc", options=[1, 5, 10, 100, 500], value=5, help="Controls how data is grouped for Topic Modeling. 1 = Line-by-line analysis (detailed). 500 = Group many lines (high-level themes).")
    st.session_state['sketch'].set_batch_size(doc_granularity)
    
    if 'last_gran' not in st.session_state: st.session_state['last_gran'] = doc_granularity
    if st.session_state['last_gran'] != doc_granularity:
        if st.session_state['sketch'].total_rows_processed > 0:
            reset_sketch()
            st.warning("Granularity changed. Data reset.")
        st.session_state['last_gran'] = doc_granularity

    topic_model_type = st.selectbox("Topic Model", ["LDA", "NMF"], help="LDA is probabilistic (good for long, mixed text). NMF is linear (good for short, distinct text like chats).")
    n_topics = st.slider("Topics", 2, 10, 4, help="The number of distinct themes the algorithm will attempt to find.")

# --- REFINERY UTILITY ---
with st.expander("🛠️ Data Refinery"):
    ref_file = st.file_uploader("CSV to Refine", type=['csv'])
    if ref_file and st.button("🚀 Run Refinery"):
        zip_data = perform_refinery_job(ref_file, 50000, clean_conf)
        if zip_data: st.download_button("Download ZIP", zip_data, "refined.zip", "application/zip")

# --- SCANNING ---
all_inputs = list(uploaded_files) if uploaded_files else []
if url_input:
    for u in url_input.split('\n'):
        if u.strip(): 
            txt = fetch_url_content(u.strip())
            if txt: 
                all_inputs.append(VirtualFile(f"url_{hash(u)}.txt", txt))
                time.sleep(URL_SCRAPE_RATE_LIMIT_SECONDS) # RATE LIMITING

if manual_input: all_inputs.append(VirtualFile("manual.txt", manual_input))

if all_inputs:
    st.subheader("🚀 Scanning Phase")
    # Config gathering loop
    file_configs = {}
    
    for idx, f in enumerate(all_inputs):
        # SAFETY: Wrap each file in try/except so one bad file doesn't crash the whole UI
        try:
            # SECURITY: RESOURCE LIMIT CHECK
            if f.getbuffer().nbytes > MAX_FILE_SIZE_MB * 1024 * 1024:
                st.error(f"❌ File **{f.name}** exceeds {MAX_FILE_SIZE_MB}MB limit.")
                continue

            file_bytes, fname, lower = f.getvalue(), f.name, f.name.lower()
            
            is_csv = lower.endswith(".csv")
            is_xlsx = lower.endswith((".xlsx", ".xlsm"))
            is_json = lower.endswith(".json")
            is_vtt = lower.endswith(".vtt")
            
            # Default config
            config = {
                "read_mode": "raw lines", "delimiter": ",", "has_header": False,
                "sheet_name": None, "json_key": None, "selected_cols": [], "join_with": " "
            }
            
            # --- input options w/ data preview-
            with st.expander(f"🧩 Scan Configuration: {fname}", expanded=True):
                if is_vtt: st.info("VTT transcript detected.")
                elif is_csv:
                    try: inferred_cols = detect_csv_num_cols(file_bytes, "auto", delimiter=",")
                    except Exception: inferred_cols = 1
                    default_mode = "csv columns" if inferred_cols > 1 else "raw lines"
                    
                    config["read_mode"] = st.radio("read mode", ["raw lines", "csv columns"], index=0 if default_mode=="raw lines" else 1, key=f"csv_mode_{idx}")
                    delim_choice = st.selectbox("delimiter", [",", "tab", ";", "|"], 0, key=f"csv_delim_{idx}")
                    config["delimiter"] = {",": ",", "tab": "\t", ";": ";", "|": "|"}[delim_choice]
                    config["has_header"] = st.checkbox("header row", value=True if inferred_cols > 1 else False, key=f"csv_header_{idx}")
                    
                    if config["read_mode"] == "csv columns":
                        st.caption("🔍 Data Preview")
                        df_prev = get_csv_preview(file_bytes, "auto", config["delimiter"], config["has_header"])
                        st.dataframe(df_prev, use_container_width=True, height=150)
                        if not df_prev.empty:
                            col_names = list(df_prev.columns)
                            config["selected_cols"] = st.multiselect("Select Text Columns to Scan", col_names, [col_names[0]], key=f"csv_cols_{idx}")
                            config["join_with"] = st.text_input("join with", " ", key=f"csv_join_{idx}")
                elif is_xlsx:
                    if openpyxl:
                        sheets = get_excel_sheetnames(file_bytes)
                        config["sheet_name"] = st.selectbox("sheet", sheets or ["(none)"], 0, key=f"xlsx_sheet_{idx}")
                        config["has_header"] = st.checkbox("header row", True, key=f"xlsx_header_{idx}")
                        if config["sheet_name"]:
                            st.caption("🔍 Data Preview")
                            df_prev = get_excel_preview(file_bytes, config["sheet_name"], config["has_header"])
                            st.dataframe(df_prev, use_container_width=True, height=150)
                            if not df_prev.empty:
                                col_names = list(df_prev.columns)
                                config["selected_cols"] = st.multiselect("Select Text Columns to Scan", col_names, [col_names[0]], key=f"xlsx_cols_{idx}")
                                config["join_with"] = st.text_input("join with", " ", key=f"xlsx_join_{idx}")

                elif is_json:
                    st.info("JSON/JSONL File.")
                    config["json_key"] = st.text_input("Key to Extract", "", key=f"json_key_{idx}")
            
            file_configs[idx] = config

            if st.button(f"⚡ Start Scan/RE-Scan: {fname}", key=f"btn_scan_{idx}"):
                if clear_on_scan: reset_sketch()
                bar = st.progress(0)
                status = st.empty()
                c = file_configs[idx]
                
                # Setup Iterator
                rows_iter = iter([])
                approx = 0
                
                if is_vtt:
                    rows_iter = read_rows_vtt(file_bytes, "auto")
                    approx = estimate_row_count_from_bytes(file_bytes)
                elif is_json:
                    rows_iter = read_rows_json(file_bytes, c["json_key"] if c["json_key"] else None)
                    approx = estimate_row_count_from_bytes(file_bytes)
                elif is_csv:
                    if c["read_mode"] == "raw lines":
                        rows_iter = read_rows_raw_lines(file_bytes, "auto")
                    else:
                        rows_iter = iter_csv_selected_columns(file_bytes, "auto", c["delimiter"], c["has_header"], c["selected_cols"], c["join_with"])
                    approx = estimate_row_count_from_bytes(file_bytes)
                elif is_xlsx and openpyxl:
                    if c["sheet_name"]:
                        rows_iter = iter_excel_selected_columns(file_bytes, c["sheet_name"], c["has_header"], c["selected_cols"], c["join_with"])
                        approx = excel_estimate_rows(file_bytes, c["sheet_name"], c["has_header"])
                elif lower.endswith(".pdf"):
                    rows_iter = read_rows_pdf(file_bytes)
                elif lower.endswith(".pptx"):
                    rows_iter = read_rows_pptx(file_bytes)
                else:
                    rows_iter = read_rows_raw_lines(file_bytes)
                    approx = estimate_row_count_from_bytes(file_bytes)
                
                def mk_cb(total):
                    def cb(n): 
                        if total: 
                            bar.progress(min(99, int(n*100/total)))
                            status.text(f"Processed {n:,} / {total:,} rows")
                        else:
                            status.text(f"Processed {n:,} rows")
                    return cb
                
                stats = Counter()
                process_chunk_iter(rows_iter, clean_conf, proc_conf, st.session_state['sketch'], mk_cb(approx), stats)
                
                bar.progress(100)
                
                # UX Feedback for Empty Scans
                if sum(stats.values()) == 0:
                    status.warning("⚠️ Scan finished, but 0 valid tokens were found.")
                    st.markdown("""
                    **Possible Reasons:**
                    *   **Numeric Data:** Did you scan a column of ID numbers while 'Drop Integers' is checked?
                    *   **Strict Filters:** Are words shorter than 'Min Word Len'?
                    """)
                else:
                    status.success(f"Scan Complete! Captured {sum(stats.values()):,} tokens.")
                
                # Quick View
                if stats:
                    st.markdown("##### 📄 Quick View: This File")
                    color_func = None
                    if enable_sentiment:
                        top_keys = [k for k,v in stats.most_common(1000)]
                        sentiments = get_sentiments(analyzer, tuple(top_keys))
                        color_func = create_sentiment_color_func(sentiments, pos_color, neg_color, neu_color, pos_threshold, neg_threshold)
                    
                    fig, _ = build_wordcloud_figure_from_counts(stats, max_words, width, height, bg_color, colormap, combined_font_path, random_state, color_func)
                    st.pyplot(fig)
                    plt.close(fig)
                
                if not clear_on_scan: st.rerun()
        
        except Exception as e:
            st.error(f"❌ Error rendering UI for file **{f.name}**: {str(e)}")

# ----------------------------
# ANALYSIS PHASE (Reads from Sketch)
# ---------------------------
scanner = st.session_state['sketch']
combined_counts = scanner.global_counts
combined_bigrams = scanner.global_bigrams

if combined_counts:
    st.divider()
    st.header("📊 Analysis Phase")
    render_analyst_help() 
    
    st.download_button(
        label="💾 Download Sketch (.json)",
        data=scanner.to_json(),
        file_name="data_sketch.json",
        mime="application/json"
    )
    
    st.info(f"Analyzing Sketch of {scanner.total_rows_processed:,} total rows.")
    
    # Sentiment Calculation
    term_sentiments = {}
    if enable_sentiment:
        top_keys = [k for k,v in combined_counts.most_common(SENTIMENT_ANALYSIS_TOP_N)]
        term_sentiments = get_sentiments(analyzer, tuple(top_keys))
        if proc_conf.compute_bigrams:
            top_bg_keys = [" ".join(k) for k,v in combined_bigrams.most_common(2000)]
            term_sentiments.update(get_sentiments(analyzer, tuple(top_bg_keys)))

    st.subheader("🔍 Bayesian Theme Discovery")
    
    # LOGIC FIX: Warn about truncated topic models
    if scanner.limit_reached:
        st.warning(f"⚠️ **Topic Model Limit Reached:** The analysis used the first {MAX_TOPIC_DOCS:,} document samples. Global word counts are accurate, but topics may not reflect the entire dataset.")

    with st.expander(f"🤔 How this works ({topic_model_type}) & Troubleshooting", expanded=False):
        n_docs = len(scanner.topic_docs)
        st.markdown(f"**Analysis Basis:** The model is learning from **{n_docs} synthetic documents** generated during the scan.")
        if n_docs < 10:
            st.warning("⚠️ **Low Resolution:** Very few documents. Decrease 'Rows per Document' to 1 for better results.")
        if topic_model_type == "LDA":
            st.markdown("**LDA** assumes documents are mixtures of topics. Best for long text.")
        else:
            st.markdown("**NMF** assumes documents belong to distinct topics. Best for short text.")

    if len(scanner.topic_docs) > 0 and DictVectorizer:
        with st.spinner(f"Running {topic_model_type} Topic Modeling..."):
            topics = perform_topic_modeling(scanner.topic_docs, n_topics, topic_model_type)
        
        if topics:
            cols = st.columns(len(topics))
            for idx, topic in enumerate(topics):
                with cols[idx]:
                    st.markdown(f"**Topic {topic['id']}**")
                    for w in topic['words']:
                        st.markdown(f"`{w}`")
        else:
            st.warning("Not enough distinct data to detect topics.")
    
    st.divider()

    st.subheader("🖼️ Combined Word Cloud")
    try:
        c_color_func = None
        if enable_sentiment: c_color_func = create_sentiment_color_func(term_sentiments, pos_color, neg_color, neu_color, pos_threshold, neg_threshold)
        fig, _ = build_wordcloud_figure_from_counts(combined_counts, max_words, width, height, bg_color, colormap, combined_font_path, random_state, c_color_func)
        st.pyplot(fig, use_container_width=True)
        st.download_button("📥 download combined png", fig_to_png_bytes(fig), "combined_wc.png", "image/png")
        plt.close(fig); gc.collect()
    except MemoryError: st.error("memory error: reduce image size.")

    st.divider()
    
    text_stats = calculate_text_stats(combined_counts, scanner.total_rows_processed)

    show_graph = proc_conf.compute_bigrams and combined_bigrams and st.checkbox("🕸️ Show Network Graph & Advanced Analytics (uncheck then re-check if graph is blank)", value=True)
    
    # --- Bayesian Sentiment Inference

    if enable_sentiment and beta_dist:
        st.subheader("⚖️ Bayesian Sentiment Inference")
        
        # <--- INSERT THIS BLOCK --->
        with st.expander("🧠 How to read this chart (and why it matters)", expanded=False):
            st.markdown("""
            **The Problem:** Standard sentiment analysis gives you a single number (e.g., "52% Positive"). But is that 52% based on 5 tweets or 5 million? A single number hides that uncertainty.
            
            **The Solution:** This chart calculates the **Probability** of the true sentiment.
            *   **The Curve (PDF):** Represents likelihood. The higher the peak, the more likely that specific sentiment score is the "truth."
            *   **The Shape:** 
                *   **Narrow & Tall:** We have lots of data. We are highly confident the sentiment is exactly here.
                *   **Wide & Flat:** We don't have enough data. The true sentiment could be almost anything.
            *   **The Green Zone (95% CI):** There is a 95% probability the "True" sentiment falls within this range. 
            
            **Decision Tip:** If the green zone is very wide (e.g., spanning 30% to 70%), **do not** make business decisions based on sentiment yet; you need more data.
            """)
        

        bayes_result = perform_bayesian_sentiment_analysis(combined_counts, term_sentiments, pos_threshold, neg_threshold) 
        if bayes_result:
            b_col1, b_col2 = st.columns([1, 2])
            with b_col1:
                st.metric("Positive Words Observed", f"{bayes_result['pos_count']:,}")
                st.metric("Negative Words Observed", f"{bayes_result['neg_count']:,}")
                st.info(f"Mean Expected Positive Rate: **{bayes_result['mean_prob']:.1%}**")
                st.success(f"95% Credible Interval:\n**{bayes_result['ci_low']:.1%} — {bayes_result['ci_high']:.1%}**")
            with b_col2:
                fig_bayes, ax_bayes = plt.subplots(figsize=(8, 4))
                ax_bayes.plot(bayes_result['x_axis'], bayes_result['pdf_y'], lw=2, color='blue', label='Posterior PDF')
                ax_bayes.fill_between(bayes_result['x_axis'], 0, bayes_result['pdf_y'], 
                                    where=(bayes_result['x_axis'] > bayes_result['ci_low']) & (bayes_result['x_axis'] < bayes_result['ci_high']),
                                    color='green', alpha=0.3, label='95% Credible Interval')
                ax_bayes.set_title("Bayesian Update of Sentiment Confidence", fontsize=10)
                ax_bayes.legend()
                ax_bayes.grid(True, alpha=0.2)
                st.pyplot(fig_bayes)
                plt.close(fig_bayes)
        st.divider()

    if show_graph:
        st.subheader("🔗 Network Graph")
        with st.expander("🛠️ Graph Settings & Physics", expanded=False):
            c1, c2, c3 = st.columns(3)
            min_edge_weight = c1.slider("Min Link Frequency", 2, 100, 2)
            max_nodes_graph = c1.slider("Max Nodes", 10, 200, 80)
            repulsion_val = c2.slider("Repulsion", 100, 3000, 1000)
            edge_len_val = c2.slider("Edge Length", 50, 500, 250)
            physics_enabled = c3.checkbox("Enable Physics", True)
            directed_graph = c3.checkbox("Directed Arrows", False)
            color_mode = c3.radio("Color By:", ["Community (Topic)", "Sentiment"], index=0)

        G = nx.DiGraph() if directed_graph else nx.Graph()
        filtered_bigrams = {k: v for k, v in combined_bigrams.items() if v >= min_edge_weight}
        sorted_connections = sorted(filtered_bigrams.items(), key=lambda x: x[1], reverse=True)[:max_nodes_graph]
        
        if sorted_connections:
            G.add_edges_from((src, tgt, {'weight': w}) for (src, tgt), w in sorted_connections)
            
            try: deg_centrality = nx.degree_centrality(G)
            except: deg_centrality = {n: 1 for n in G.nodes()}
            community_map = {}
            ai_cluster_info = "" 
            
            if color_mode == "Community (Topic)":
                G_undir = G.to_undirected() if directed_graph else G
                try:
                    communities = nx_comm.greedy_modularity_communities(G_undir)
                    cluster_descriptions = []
                    for group_id, community in enumerate(communities):
                        top_in_cluster = sorted(list(community), key=lambda x: combined_counts[x], reverse=True)[:5]
                        cluster_descriptions.append(f"- Cluster {group_id+1}: {', '.join(top_in_cluster)}")
                        for node in community: community_map[node] = group_id
                    ai_cluster_info = "\n".join(cluster_descriptions)
                except Exception as e: pass

            community_colors = ["#FF4B4B", "#4589ff", "#ffa421", "#3cdb82", "#8b46ff", "#ff4b9f", "#00c0f2"]
            nodes, edges = [], []
            for node_id in G.nodes():
                size = 15 + (deg_centrality.get(node_id, 0) * 80)
                if color_mode == "Sentiment":
                    node_color = neu_color
                    if enable_sentiment:
                        score = term_sentiments.get(node_id, 0)
                        if score >= pos_threshold: node_color = pos_color
                        elif score <= neg_threshold: node_color = neg_color
                else:
                    group_id = community_map.get(node_id, 0)
                    node_color = community_colors[group_id % len(community_colors)]

                nodes.append(Node(id=node_id, label=node_id, size=size, color=node_color, title=f"Term: {node_id}\nFreq: {combined_counts.get(node_id, 0)}", font={'color': 'white', 'size': 20, 'strokeWidth': 4, 'strokeColor': '#000000'}))

            for (source, target), weight in sorted_connections:
                width = 1 + math.log(weight) * 0.8
                edges.append(Edge(source=source, target=target, width=width, color="#e0e0e0"))
            
            config = Config(width=1000, height=700, directed=directed_graph, physics=physics_enabled, hierarchy=False, interaction={"navigationButtons": True, "zoomView": True}, physicsSettings={"solver": "forceAtlas2Based", "forceAtlas2Based": {"gravitationalConstant": -abs(repulsion_val), "springLength": edge_len_val, "springConstant": 0.05, "damping": 0.4}})
            st.info("💡 **Navigation Tip:** Use the buttons in the **bottom-right** of the graph to Zoom & Pan.")
            agraph(nodes=nodes, edges=edges, config=config)

            st.markdown("### 📊 Graph Analytics")
            tab1, tab2, tab3, tab4 = st.tabs(["Basic Stats", "Top Nodes", "Text Stats", "🔥 Heatmap"])
            with tab1:
                col_b1, col_b2, col_b3 = st.columns(3)
                col_b1.metric("Nodes", G.number_of_nodes())
                col_b2.metric("Edges", G.number_of_edges())
                try: col_b3.metric("Density", f"{nx.density(G):.4f}")
                except: pass
            with tab2:
                node_weights = {n: 0 for n in G.nodes()}
                for u, v, data in G.edges(data=True):
                    w = data.get('weight', 1)
                    node_weights[u] += w
                    node_weights[v] += w
                st.dataframe(pd.DataFrame(list(node_weights.items()), columns=["Node", "Weighted Degree"]).sort_values("Weighted Degree", ascending=False).head(50), use_container_width=True)
            with tab3:
                col_s1, col_s2, col_s3, col_s4 = st.columns(4)
                col_s1.metric("Total Tokens", f"{text_stats['Total Tokens']:,}")
                col_s2.metric("Unique Vocab", f"{text_stats['Unique Vocabulary']:,}")
                col_s3.metric("Lexical Diversity", f"{text_stats['Lexical Diversity']}")
                col_s4.metric("Avg Word Len", f"{text_stats['Avg Word Length']}")
            
            # --- HEATMAP VISUALIZATION ---
            with tab4:
                st.caption("Shows how often the Top 20 terms appear next to each other.")
                top_20 = [w for w, c in combined_counts.most_common(20)]
                mat_size = len(top_20)
                if mat_size > 1:
                    heatmap_matrix = np.zeros((mat_size, mat_size))
                    for i, w1 in enumerate(top_20):
                        for j, w2 in enumerate(top_20):
                            if i != j:
                                val = combined_bigrams.get((w1, w2), 0) + combined_bigrams.get((w2, w1), 0)
                                heatmap_matrix[i][j] = val

                    fig_h, ax_h = plt.subplots(figsize=(10, 8))
                    im = ax_h.imshow(heatmap_matrix, cmap="Blues")
                    ax_h.set_xticks(np.arange(mat_size))
                    ax_h.set_yticks(np.arange(mat_size))
                    ax_h.set_xticklabels(top_20, rotation=45, ha="right")
                    ax_h.set_yticklabels(top_20)
                    
                    for i in range(mat_size):
                        for j in range(mat_size):
                            if heatmap_matrix[i, j] > 0:
                                ax_h.text(j, i, int(heatmap_matrix[i, j]), ha="center", va="center", color="black" if heatmap_matrix[i,j] < heatmap_matrix.max()/2 else "white", fontsize=8)
                    
                    ax_h.set_title("Top 20 Terms Co-occurrence")
                    plt.tight_layout()
                    st.pyplot(fig_h)
                    plt.close(fig_h)
                else:
                    st.info("Not enough data for heatmap.")

    else:
        st.subheader("📈 Text Statistics")
        col_s1, col_s2, col_s3, col_s4 = st.columns(4)
        col_s1.metric("Total Tokens", f"{text_stats['Total Tokens']:,}")
        col_s2.metric("Unique Vocab", f"{text_stats['Unique Vocabulary']:,}")
        col_s3.metric("Lexical Diversity", f"{text_stats['Lexical Diversity']}")
        col_s4.metric("Avg Word Len", f"{text_stats['Avg Word Length']}")

else: 
    st.info("Uploaded data will be processed into the sketch above.")


# --- ai analytics section
if combined_counts and st.session_state['authenticated']:
    st.divider()
    st.subheader("🤖 AI Analyst")
    
    # Context Preparation
    top_unigrams = [w for w, c in combined_counts.most_common(100)]
    top_bigrams = [" ".join(bg) for bg, c in combined_bigrams.most_common(30)] if proc_conf.compute_bigrams else []
    g_context = locals().get("ai_cluster_info", "(Graph clustering not run)")
    
    ai_context_str = f"""
    Top 100 Unigrams: {', '.join(top_unigrams)}
    Top 30 Bigrams: {', '.join(top_bigrams)}
    Graph Clusters: {g_context}
    """

    col_ai_1, col_ai_2 = st.columns(2)
    
    # 1. Automatic Analysis
    with col_ai_1:
        st.markdown("**1. One-Click Theme Detection**")
        if st.button("✨ Identify Key Themes", type="primary"):
            with st.status("Analyzing...", expanded=True):
                system_prompt = "You are a qualitative data analyst. Analyze the provided word frequency lists to identify 3 key themes, potential anomalies, and a summary of the subject matter."
                user_prompt = f"Data Context:\n{ai_context_str}"
                response = call_llm_and_track_cost(system_prompt, user_prompt, ai_config)
                st.session_state["ai_response"] = response
                st.rerun()

    # 2. Free Form Question
    with col_ai_2:
        st.markdown("**2. Ask the Data**")
        user_question = st.text_area("Ask a specific question:", height=100, placeholder="e.g., 'What are the main complaints about pricing?'")
        if st.button("Ask Question"):
            if user_question.strip():
                with st.status("Thinking...", expanded=True):
                    system_prompt = "You are an expert analyst. Answer the user's question based ONLY on the provided summary statistics (word counts and associations). If you cannot answer from the data, say so."
                    user_prompt = f"Data Context:\n{ai_context_str}\n\nUser Question: {user_question}"
                    response = call_llm_and_track_cost(system_prompt, user_prompt, ai_config)
                    st.session_state["ai_response"] = f"**Q: {user_question}**\n\n{response}"
                    st.rerun()
            else:
                st.warning("Please enter a question.")

    if st.session_state.get("ai_response"):
        st.divider()
        st.markdown("### 📋 AI Output")
        st.markdown(st.session_state["ai_response"])

# ---tables
if combined_counts:
    st.divider()
    st.subheader(f"📊 Frequency Tables (Top {top_n})")
    most_common = combined_counts.most_common(top_n)
    
    data = []
    if enable_sentiment:
        for w, f in most_common:
            score = term_sentiments.get(w, 0.0)
            category = get_sentiment_category(score, pos_threshold, neg_threshold)
            data.append([w, f, score, category])
    else:
        data = [[w, f] for w, f in most_common]

    cols = ["word", "count"] + (["sentiment", "category"] if enable_sentiment else [])
    st.dataframe(pd.DataFrame(data, columns=cols), use_container_width=True)
    
    if proc_conf.compute_bigrams and combined_bigrams:
        st.write("Bigrams (By Frequency)")
        top_bg = combined_bigrams.most_common(top_n)
        
        bg_data = []
        if enable_sentiment:
            for bg_tuple, f in top_bg:
                bg_str = " ".join(bg_tuple)
                score = term_sentiments.get(bg_str, 0.0)
                category = get_sentiment_category(score, pos_threshold, neg_threshold)
                bg_data.append([bg_str, f, score, category])
        else:
            bg_data = [[" ".join(bg), f] for bg, f in top_bg]

        bg_cols = ["bigram", "count"] + (["sentiment", "category"] if enable_sentiment else [])
        st.dataframe(pd.DataFrame(bg_data, columns=bg_cols), use_container_width=True)

        with st.expander("🔬 Phrase Significance (NPMI Score)", expanded=False):
            st.markdown("""
            **NPMI (Normalized Pointwise Mutual Information)** finds words that *belong* together, rather than just words that appear often.
            *   High Score (> 0.5): Strong association (e.g., "Artificial Intelligence").
            *   Low Score (< 0.1): Random association (e.g., "of the").
            """)
            df_npmi = calculate_npmi(combined_bigrams, combined_counts, scanner.total_rows_processed, min_freq=3)
            st.dataframe(df_npmi.head(top_n), use_container_width=True)

st.markdown("---")
st.markdown(
    "<div style='text-align: center; color: #808080; font-size: 12px;'>"
    "Open Source software licensed under the MIT License."
    "</div>", 
    unsafe_allow_html=True
)
