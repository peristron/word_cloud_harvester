import re
import html
import string
import csv
import io
import json
import pandas as pd
import time
from collections import Counter
from typing import List, Tuple, Iterable, Dict, Optional, Callable
from itertools import pairwise
from wordcloud import STOPWORDS
import openpyxl

# Patterns
HTML_TAG_RE = re.compile(r"<[^>]+>")
CHAT_ARTIFACT_RE = re.compile(
    r":\w+:"
    r"|\b(?:monday|tuesday|wednesday|thursday|friday|saturday|sunday|today|yesterday) at \d{1,2}:\d{2}\b"
    r"|\b\d+\s+repl(?:y|ies)\b"
    r"|\d{2}:\d{2}:\d{2}\.\d{3}\s+-->\s+\d{2}:\d{2}:\d{2}\.\d{3}"
    r"|\[[^\]]+\]",
    flags=re.IGNORECASE
)

def default_prepositions() -> set:
    return {'about', 'above', 'across', 'after', 'against', 'along', 'among', 'around', 'at', 'before', 'behind', 'below', 'beneath', 'beside', 'between', 'beyond', 'but', 'by', 'concerning', 'despite', 'down', 'during', 'except', 'for', 'from', 'in', 'inside', 'into', 'like', 'near', 'of', 'off', 'on', 'onto', 'out', 'outside', 'over', 'past', 'regarding', 'since', 'through', 'throughout', 'to', 'toward', 'under', 'underneath', 'until', 'up', 'upon', 'with', 'within', 'without'}

def build_punct_translation(keep_hyphens: bool, keep_apostrophes: bool) -> dict:
    punct = string.punctuation
    if keep_hyphens: punct = punct.replace("-", "")
    if keep_apostrophes: punct = punct.replace("'", "")
    return str.maketrans("", "", punct)

def build_phrase_pattern(phrases: List[str]) -> Optional[re.Pattern]:
    if not phrases: return None
    escaped = [re.escape(p) for p in phrases if p]
    if not escaped: return None
    return re.compile(rf"\b(?:{'|'.join(escaped)})\b", flags=re.IGNORECASE)

# --- THIS IS THE FUNCTION THAT WAS MISSING ---
def parse_user_stopwords(raw: str) -> Tuple[List[str], List[str]]:
    raw = raw.replace("\n", ",").replace(".", ",")
    phrases, singles = [], []
    for item in [x.strip() for x in raw.split(",") if x.strip()]:
        if " " in item: phrases.append(item.lower())
        else: singles.append(item.lower())
    return phrases, singles

def is_url_token(tok: str) -> bool:
    t = tok.strip("()[]{}<>,.;:'\"!?").lower()
    if not t: return False
    return ("://" in t) or t.startswith("www.")

def clean_and_tokenize(
    text: str,
    remove_chat: bool, remove_html: bool, unescape: bool, remove_urls: bool,
    trans_map: dict,
    stopwords: set,
    phrase_pattern: Optional[re.Pattern],
    min_len: int, drop_int: bool
) -> List[str]:
    """Core stateless tokenizer used by both App and Harvester."""
    if not text: return []
    if remove_chat: text = CHAT_ARTIFACT_RE.sub(" ", text)
    if remove_html: text = HTML_TAG_RE.sub(" ", text)
    if unescape:
        try: text = html.unescape(text)
        except: pass
    text = text.lower()
    if phrase_pattern: text = phrase_pattern.sub(" ", text)
    
    tokens = []
    for t in text.split():
        if remove_urls and is_url_token(t): continue
        t = t.translate(trans_map)
        if not t or len(t) < min_len or (drop_int and t.isdigit()) or t in stopwords: continue
        tokens.append(t)
    return tokens

# --- Readers ---

def read_rows_raw_lines(file_bytes: bytes, encoding_choice: str = "auto") -> Iterable[str]:
    def _iter_with_encoding(enc: str):
        bio = io.BytesIO(file_bytes)
        with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline=None) as wrapper:
            for line in wrapper: yield line.rstrip("\r\n")
    if encoding_choice == "latin-1": yield from _iter_with_encoding("latin-1")
    else: yield from _iter_with_encoding("utf-8")

def read_rows_vtt(file_bytes: bytes, encoding_choice: str = "auto") -> Iterable[str]:
    for line in read_rows_raw_lines(file_bytes, encoding_choice):
        line = line.strip()
        if not line or line == "WEBVTT" or "-->" in line or line.isdigit(): continue
        if ":" in line:
            parts = line.split(":", 1)
            if len(parts) > 1 and len(parts[0]) < 30 and " " in parts[0]:
                yield parts[1].strip()
                continue
        yield line

def read_rows_pdf(file_bytes: bytes) -> Iterable[str]:
    import pypdf
    bio = io.BytesIO(file_bytes)
    try:
        reader = pypdf.PdfReader(bio)
        for page in reader.pages:
            text = page.extract_text()
            if text: yield text
    except Exception: yield ""

def read_rows_pptx(file_bytes: bytes) -> Iterable[str]:
    import pptx
    bio = io.BytesIO(file_bytes)
    try:
        prs = pptx.Presentation(bio)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    if shape.text: yield shape.text
    except Exception: yield ""

def get_csv_preview(file_bytes: bytes, encoding_choice: str = "auto") -> pd.DataFrame:
    enc = "latin-1" if encoding_choice == "latin-1" else "utf-8"
    bio = io.BytesIO(file_bytes)
    try:
        df = pd.read_csv(bio, header=0, nrows=5, encoding=enc, on_bad_lines='skip')
        return df
    except: return pd.DataFrame()

def iter_csv_selected_columns(file_bytes: bytes, encoding_choice: str, delimiter: str, has_header: bool, selected_columns: List[str]) -> Iterable[str]:
    enc = "latin-1" if encoding_choice == "latin-1" else "utf-8"
    bio = io.BytesIO(file_bytes)
    with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline="") as wrapper:
        rdr = csv.reader(wrapper, delimiter=delimiter)
        first = next(rdr, None)
        if first is None: return
        
        name_to_idx = {n: i for i, n in enumerate(first)}
        idxs = [name_to_idx[n] for n in selected_columns if n in name_to_idx]

        for row in rdr:
            vals = [row[i] if i < len(row) else "" for i in idxs]
            if any(vals): yield " ".join(str(v) for v in vals if v)

def perform_topic_modeling(file_counts: List[Counter], n_topics: int = 4, top_n_words: int = 6, model_type: str = "LDA") -> Optional[List[Dict]]:
    from sklearn.decomposition import LatentDirichletAllocation, NMF
    from sklearn.feature_extraction import DictVectorizer
    
    valid_counts = [c for c in file_counts if c and len(c) > 0]
    if not valid_counts: return None
    
    vectorizer = DictVectorizer(sparse=True)
    dtm = vectorizer.fit_transform(valid_counts)
    
    model = None
    if model_type == "LDA":
        model = LatentDirichletAllocation(n_components=n_topics, random_state=42, learning_method='batch', max_iter=10)
    elif model_type == "NMF":
        model = NMF(n_components=n_topics, random_state=42, init='nndsvd')
    
    if not model: return None
    model.fit(dtm)
    
    feature_names = vectorizer.get_feature_names_out()
    topics = []
    for topic_idx, topic in enumerate(model.components_):
        top_indices = topic.argsort()[:-top_n_words - 1:-1]
        topics.append({"id": topic_idx + 1, "words": [feature_names[i] for i in top_indices]})
    return topics
